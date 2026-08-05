from __future__ import annotations

import asyncio
import inspect
import json
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path
from typing import Any
from uuid import UUID

from pypdf import PdfReader

from cognigraph.domain.base import json_compatible
from cognigraph.domain.documents import BoundingBox
from cognigraph.ingestion.models import ParsedBlock, ParsedDocument
from cognigraph.ingestion.ocr_adapter import OCRUnavailableError, PaddleOCRAdapter
from cognigraph.ingestion.vision_adapter import VisionParser, VisionParserError


class DocumentParser:
    """Docling-first parser with deterministic OCR and vision fallbacks.

    CPU-bound Docling/OCR work runs in a worker thread through ``parse_async``;
    provider calls stay on the application's event loop so async database audit
    sinks are never reused across event loops.
    """

    def __init__(
        self,
        ocr: PaddleOCRAdapter | None = None,
        vision: VisionParser | None = None,
        *,
        min_text_quality: float = 0.2,
        ocr_enabled: bool = True,
    ) -> None:
        if not 0.0 <= min_text_quality <= 1.0:
            raise ValueError("min_text_quality must be between 0 and 1")
        self.ocr = ocr or PaddleOCRAdapter()
        self.vision = vision
        self.min_text_quality = min_text_quality
        self.ocr_enabled = ocr_enabled

    def parse(
        self,
        path: Path,
        mime_type: str,
        *,
        workspace_id: UUID | None = None,
        document_id: UUID | None = None,
    ) -> ParsedDocument:
        """Synchronous compatibility entry point for scripts and unit tests.

        Async applications must use ``parse_async``.  Refusing to bridge a
        running event loop prevents asyncpg/SQLAlchemy objects from crossing
        loop boundaries during Vision model auditing.
        """

        parsed = self._parse_without_vision(path, mime_type)
        if not self._requires_vision(parsed, path, mime_type):
            return parsed
        try:
            asyncio.get_running_loop()
        except RuntimeError:
            return asyncio.run(
                self._augment_with_vision(
                    parsed,
                    path,
                    mime_type,
                    workspace_id=workspace_id,
                    document_id=document_id,
                )
            )
        raise RuntimeError("DocumentParser.parse cannot run Vision inside an active event loop")

    def _parse_without_vision(self, path: Path, mime_type: str) -> ParsedDocument:
        docling_warning: str | None = None
        try:
            parsed = self._parse_docling(path)
            if parsed.blocks and _text_quality(parsed.blocks) >= self.min_text_quality:
                # Docling can successfully extract the text layer from one page
                # of a mixed PDF while silently leaving scanned pages empty.
                # Complete those pages before treating the document as parsed;
                # otherwise the OCR/Vision fallback is never reached.
                if _is_pdf_input(path, mime_type):
                    parsed = self._complete_pdf_pages(parsed, path)
                return _ensure_parser_metadata(parsed, "docling")
            docling_warning = "Docling produced no usable text blocks"
            if parsed.blocks:
                docling_warning = "Docling text quality was below the configured threshold"
        except (ImportError, RuntimeError, ValueError, TypeError, AttributeError, OSError) as exc:
            docling_warning = f"Docling unavailable or failed: {type(exc).__name__}"

        try:
            fallback = self._parse_fallback(path, mime_type)
        except (ImportError, RuntimeError, ValueError, TypeError, AttributeError, OSError) as exc:
            fallback = ParsedDocument(
                parser_name="fallback",
                parser_version="1",
                page_count=1,
                blocks=[],
                raw_payload={},
                warnings=[f"fallback parser failed: {type(exc).__name__}"],
            )
        if docling_warning is not None:
            fallback.warnings.insert(0, docling_warning)
        visual_input = _is_visual_input(path, mime_type)
        if fallback.blocks and (
            not visual_input or _parsed_quality(fallback) >= self.min_text_quality
        ):
            return _ensure_parser_metadata(fallback, fallback.parser_name)

        return _ensure_parser_metadata(fallback, fallback.parser_name)

    def _complete_pdf_pages(self, parsed: ParsedDocument, path: Path) -> ParsedDocument:
        """OCR only PDF pages that remain unresolved after Docling extraction."""

        try:
            reader = PdfReader(str(path))
        except (OSError, ValueError, TypeError) as exc:
            return parsed.model_copy(
                update={
                    "warnings": [
                        *parsed.warnings,
                        f"PDF page coverage check failed: {type(exc).__name__}",
                    ]
                }
            )

        page_count = max(parsed.page_count, len(reader.pages))
        blocks_by_page: dict[int, list[ParsedBlock]] = {}
        for block in parsed.blocks:
            if block.page_number is not None:
                blocks_by_page.setdefault(block.page_number, []).append(block)

        unresolved_pages: list[int] = []
        text_layer_blocks: list[ParsedBlock] = []
        for page_number, page in enumerate(reader.pages, start=1):
            docling_quality = _text_quality(blocks_by_page.get(page_number, []))
            if docling_quality >= self.min_text_quality:
                continue
            page_text = (page.extract_text() or "").strip()
            extracted_quality = (
                _text_quality([ParsedBlock(text=page_text, page_number=page_number)])
                if page_text
                else 0.0
            )
            if extracted_quality >= self.min_text_quality:
                # Docling can miss an ordinary text-layer page even though
                # pypdf can extract it. Preserve that source content rather
                # than merely using it to suppress OCR for the page.
                text_layer_blocks.append(ParsedBlock(text=page_text, page_number=page_number))
            else:
                unresolved_pages.append(page_number)

        completed_blocks = _merge_block_lists(parsed.blocks, text_layer_blocks)
        parser_chain = list(parsed.parser_chain)
        if text_layer_blocks and "pypdf" not in parser_chain:
            parser_chain.append("pypdf")

        if not unresolved_pages:
            return parsed.model_copy(
                update={
                    "page_count": page_count,
                    "blocks": completed_blocks,
                    "parser_chain": parser_chain,
                }
            )
        if not self.ocr_enabled:
            return parsed.model_copy(
                update={
                    "page_count": page_count,
                    "blocks": completed_blocks,
                    "parser_chain": parser_chain,
                    "raw_payload": {
                        **parsed.raw_payload,
                        "unresolved_pages": unresolved_pages,
                    },
                    "warnings": [
                        *parsed.warnings,
                        f"PDF has {len(unresolved_pages)} page(s) requiring OCR; OCR is disabled.",
                    ],
                }
            )

        try:
            ocr = self.ocr.parse_pdf(path, page_numbers=set(unresolved_pages))
        except OCRUnavailableError as exc:
            return parsed.model_copy(
                update={
                    "page_count": page_count,
                    "blocks": completed_blocks,
                    "parser_chain": parser_chain,
                    "raw_payload": {
                        **parsed.raw_payload,
                        "unresolved_pages": unresolved_pages,
                    },
                    "warnings": [
                        *parsed.warnings,
                        f"PDF has {len(unresolved_pages)} page(s) requiring OCR.",
                        str(exc),
                    ],
                }
            )

        merged_blocks = _merge_block_lists(completed_blocks, ocr.blocks)
        resolved_pages = {
            block.page_number for block in ocr.blocks if block.page_number is not None
        }
        still_unresolved = [page for page in unresolved_pages if page not in resolved_pages]
        return parsed.model_copy(
            update={
                "page_count": max(page_count, ocr.page_count),
                "blocks": merged_blocks,
                "language": parsed.language or ocr.language,
                "detected_language": parsed.detected_language or ocr.detected_language,
                "raw_payload": {
                    **parsed.raw_payload,
                    "page_count": page_count,
                    "unresolved_pages": still_unresolved,
                    "ocr": ocr.raw_payload,
                },
                "warnings": [*parsed.warnings, *ocr.warnings],
                "parser_chain": list(dict.fromkeys([*parser_chain, *ocr.parser_chain])),
                "ocr_used": True,
                "low_confidence_blocks": [
                    *parsed.low_confidence_blocks,
                    *ocr.low_confidence_blocks,
                ],
            }
        )

    async def parse_async(
        self,
        path: Path,
        mime_type: str,
        *,
        workspace_id: UUID | None = None,
        document_id: UUID | None = None,
    ) -> ParsedDocument:
        parsed = await asyncio.to_thread(
            self._parse_without_vision,
            path,
            mime_type,
        )
        if not self._requires_vision(parsed, path, mime_type):
            return parsed
        return await self._augment_with_vision(
            parsed,
            path,
            mime_type,
            workspace_id=workspace_id,
            document_id=document_id,
        )

    def _requires_vision(
        self,
        parsed: ParsedDocument,
        path: Path,
        mime_type: str,
    ) -> bool:
        if self.vision is None or not _is_visual_input(path, mime_type):
            return False
        unresolved_pages = parsed.raw_payload.get("unresolved_pages", [])
        if isinstance(unresolved_pages, list) and unresolved_pages:
            return True
        if not parsed.blocks or _parsed_quality(parsed) < self.min_text_quality:
            return True
        complex_labels = ("chart", "diagram", "figure", "formula", "image", "picture", "table")
        return any(
            any(label in block.block_type.casefold() for label in complex_labels)
            for block in parsed.blocks
        )

    async def _augment_with_vision(
        self,
        parsed: ParsedDocument,
        path: Path,
        mime_type: str,
        *,
        workspace_id: UUID | None,
        document_id: UUID | None,
    ) -> ParsedDocument:
        try:
            visual = await self._parse_with_vision_async(
                path,
                mime_type,
                workspace_id=workspace_id,
                document_id=document_id,
                page_numbers=_unresolved_pages(parsed),
            )
        except VisionParserError as exc:
            return parsed.model_copy(update={"warnings": [*parsed.warnings, str(exc)]})
        except (ImportError, OSError, RuntimeError, ValueError) as exc:
            return parsed.model_copy(
                update={
                    "warnings": [
                        *parsed.warnings,
                        f"Vision parser failed: {type(exc).__name__}",
                    ]
                }
            )
        if not visual.blocks:
            return parsed.model_copy(
                update={"warnings": [*parsed.warnings, "Vision parser returned no usable blocks"]}
            )
        return _merge_visual_results(parsed, visual)

    @staticmethod
    def _parse_docling(path: Path) -> ParsedDocument:
        try:
            # Keep the optional Docling dependency out of the default runtime and
            # type-checker import graph while still resolving its current API at
            # execution time when the documents extra is installed.
            from importlib import import_module

            converter_module = import_module("docling.document_converter")
            document_converter = converter_module.DocumentConverter
        except ImportError as exc:
            raise ImportError("Docling optional dependency is not installed") from exc
        converter = document_converter()
        result: Any = converter.convert(path)
        document = result.document
        payload: dict[str, Any] = document.export_to_dict()
        blocks: list[ParsedBlock] = []
        texts = getattr(document, "texts", [])
        heading_path: list[str] = []
        for item in texts:
            text = str(getattr(item, "text", "")).strip()
            if not text:
                continue
            page_number: int | None = None
            provenance = getattr(item, "prov", None) or []
            if provenance:
                page_number = int(getattr(provenance[0], "page_no", 0) or 0) or None
            label = str(getattr(item, "label", "paragraph"))
            normalized_label = label.casefold()
            if "title" in normalized_label or "section_header" in normalized_label:
                heading_path = [text]
            blocks.append(
                ParsedBlock(
                    text=text,
                    page_number=page_number,
                    heading_path=list(heading_path),
                    bounding_box=_docling_bounding_box(provenance[0]) if provenance else None,
                    block_type=label,
                )
            )
        if not blocks:
            markdown = str(document.export_to_markdown()).strip()
            if markdown:
                blocks.append(ParsedBlock(text=markdown, page_number=1))
        pages = getattr(document, "pages", {})
        return ParsedDocument(
            parser_name="docling",
            parser_version=_package_version("docling"),
            page_count=max(len(pages), 1 if blocks else 0),
            blocks=blocks,
            raw_payload=json.loads(json.dumps(payload, default=str)),
            parser_chain=["docling"],
        )

    def _parse_fallback(self, path: Path, mime_type: str) -> ParsedDocument:
        suffix = path.suffix.casefold()
        if suffix == ".pdf" or mime_type == "application/pdf":
            return self._parse_pdf(path)
        if suffix == ".docx":
            return self._parse_docx(path)
        if suffix == ".pptx":
            return self._parse_pptx(path)
        if suffix in {".txt", ".md"} or mime_type.startswith("text/"):
            text = path.read_text(encoding="utf-8-sig")
            blocks = [
                ParsedBlock(text=part.strip(), page_number=1)
                for part in text.split("\n\n")
                if part.strip()
            ]
            return ParsedDocument(
                parser_name="plain-text",
                parser_version="1",
                page_count=1,
                blocks=blocks,
                raw_payload={"character_count": len(text)},
                parser_chain=["plain-text"],
            )
        if mime_type.startswith("image/"):
            if not self.ocr_enabled:
                return ParsedDocument(
                    parser_name="pillow-image",
                    parser_version=_package_version("pillow"),
                    page_count=1,
                    blocks=[],
                    raw_payload={},
                    warnings=[
                        "Image OCR is disabled by configuration.",
                        "Image requires Docling, PaddleOCR, or a configured vision model.",
                    ],
                    parser_chain=["pillow-image"],
                )
            try:
                return self.ocr.parse(path)
            except OCRUnavailableError as exc:
                return ParsedDocument(
                    parser_name="pillow-image",
                    parser_version=_package_version("pillow"),
                    page_count=1,
                    blocks=[],
                    raw_payload={},
                    warnings=[str(exc), "Image requires Docling, PaddleOCR, or a vision model."],
                    parser_chain=["pillow-image"],
                )
        raise ValueError(f"no parser for {mime_type}")

    def _parse_pdf(self, path: Path) -> ParsedDocument:
        reader = PdfReader(str(path))
        blocks: list[ParsedBlock] = []
        missing_pages: list[int] = []
        for index, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            page_block = ParsedBlock(text=text, page_number=index) if text else None
            if page_block is not None and _text_quality([page_block]) >= self.min_text_quality:
                blocks.append(page_block)
            else:
                missing_pages.append(index)
        extracted = ParsedDocument(
            parser_name="pypdf",
            parser_version=_package_version("pypdf"),
            page_count=len(reader.pages),
            blocks=blocks,
            raw_payload={
                "page_count": len(reader.pages),
                "unresolved_pages": missing_pages,
            },
            parser_chain=["pypdf"],
        )
        if not missing_pages:
            return extracted
        # Render only the pages whose text layer is absent or unusable. This
        # preserves mixed PDFs containing both digital and scanned pages.
        if not self.ocr_enabled:
            extracted.warnings.append(
                f"PDF has {len(missing_pages)} page(s) requiring OCR; OCR is disabled."
            )
            return extracted
        try:
            ocr = self.ocr.parse_pdf(path, page_numbers=set(missing_pages))
        except OCRUnavailableError as exc:
            extracted.warnings.append(f"PDF has {len(missing_pages)} page(s) requiring OCR.")
            extracted.warnings.append(str(exc))
            return extracted
        resolved_pages = {
            block.page_number for block in ocr.blocks if block.page_number is not None
        }
        unresolved_pages = [page for page in missing_pages if page not in resolved_pages]
        # Preserve the text layer and append OCR lines. Page/source provenance
        # lets downstream chunking retain deterministic reading order.
        ocr.blocks = sorted(
            [*blocks, *ocr.blocks],
            key=lambda block: (
                block.page_number or 0,
                block.bounding_box.top if block.bounding_box else 0,
            ),
        )
        ocr.page_count = max(ocr.page_count, extracted.page_count)
        ocr.parser_chain = ["pypdf", *ocr.parser_chain]
        ocr.raw_payload = {
            **ocr.raw_payload,
            "page_count": extracted.page_count,
            "unresolved_pages": json_compatible(unresolved_pages),
        }
        return ocr

    @staticmethod
    def _parse_docx(path: Path) -> ParsedDocument:
        try:
            from docx import Document as DocxDocument
        except ImportError as exc:
            raise RuntimeError("python-docx is required for DOCX fallback parsing") from exc
        document = DocxDocument(str(path))
        blocks: list[ParsedBlock] = []
        heading_path: list[str] = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = paragraph.style.name if paragraph.style is not None else ""
            if style_name.startswith("Heading"):
                heading_path = [text]
            blocks.append(ParsedBlock(text=text, page_number=1, heading_path=heading_path))
        return ParsedDocument(
            parser_name="python-docx",
            parser_version=_package_version("python-docx"),
            page_count=1,
            blocks=blocks,
            raw_payload={"paragraph_count": len(document.paragraphs)},
            parser_chain=["python-docx"],
        )

    @staticmethod
    def _parse_pptx(path: Path) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("python-pptx is required for PPTX fallback parsing") from exc
        presentation = Presentation(str(path))
        blocks: list[ParsedBlock] = []
        for page_number, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                text = str(getattr(shape, "text", "")).strip()
                if text:
                    blocks.append(ParsedBlock(text=text, page_number=page_number))
        return ParsedDocument(
            parser_name="python-pptx",
            parser_version=_package_version("python-pptx"),
            page_count=len(presentation.slides),
            blocks=blocks,
            raw_payload={"slide_count": len(presentation.slides)},
            parser_chain=["python-pptx"],
        )

    async def _parse_with_vision_async(
        self,
        path: Path,
        mime_type: str,
        *,
        workspace_id: UUID | None,
        document_id: UUID | None,
        page_numbers: set[int] | None,
    ) -> ParsedDocument:
        if self.vision is None:
            raise VisionParserError("vision parser is not configured")
        resolved_workspace = workspace_id or UUID(int=0)
        resolved_document = document_id or UUID(int=0)
        if path.suffix.casefold() != ".pdf" and mime_type != "application/pdf":
            return await self.vision.parse_image(
                path,
                workspace_id=resolved_workspace,
                document_id=resolved_document,
            )

        # Render each scan page and ask the vision parser independently so page
        # numbers and bounding boxes remain auditable.
        renderer = getattr(self.ocr, "_render_pdf_pages", None)
        if not callable(renderer):
            raise VisionParserError("scan-PDF vision fallback requires a PDF renderer")
        rendered_pages = await asyncio.to_thread(
            _render_selected_pages,
            renderer,
            path,
            page_numbers,
        )
        blocks: list[ParsedBlock] = []
        warnings: list[str] = []
        page_count = 0
        try:
            for page_number, image_path in rendered_pages:
                page_count = max(page_count, page_number)
                try:
                    page = await self.vision.parse_image(
                        image_path,
                        workspace_id=resolved_workspace,
                        document_id=resolved_document,
                    )
                    for block in page.blocks:
                        blocks.append(
                            # Each request contains exactly one rendered page;
                            # the renderer page is authoritative provenance.
                            block.model_copy(update={"page_number": page_number})
                        )
                    warnings.extend(page.warnings)
                finally:
                    image_path.unlink(missing_ok=True)
        finally:
            _cleanup_render_directories(rendered_pages)
        return ParsedDocument(
            parser_name="litellm-vision",
            parser_version="1",
            page_count=page_count,
            blocks=blocks,
            language=_detect_language(blocks),
            detected_language=_detect_language(blocks),
            raw_payload={"page_count": page_count, "block_count": len(blocks)},
            warnings=warnings,
            parser_chain=["pdf-renderer", "litellm-vision"],
            vision_used=True,
            low_confidence_blocks=[
                {
                    "page_number": block.page_number,
                    "text": block.text[:200],
                    "confidence": block.confidence,
                }
                for block in blocks
                if block.confidence is not None and block.confidence < 0.6
            ],
        )


def _render_selected_pages(
    renderer: Any,
    path: Path,
    page_numbers: set[int] | None,
) -> list[tuple[int, Path]]:
    try:
        supports_selection = "page_numbers" in inspect.signature(renderer).parameters
    except (TypeError, ValueError):
        supports_selection = False
    rendered = renderer(path, page_numbers=page_numbers) if supports_selection else renderer(path)
    pages = list(rendered)
    if page_numbers is not None and not supports_selection:
        pages = [item for item in pages if item[0] in page_numbers]
    return pages


def _cleanup_render_directories(rendered_pages: list[tuple[int, Path]]) -> None:
    parents = {image_path.parent for _, image_path in rendered_pages}
    for parent in parents:
        if parent.name.startswith("cognigraph-pdf-") and parent.exists():
            try:
                parent.rmdir()
            except OSError:
                # A renderer may retain an auxiliary file. It owns cleanup for
                # anything other than the page images already removed above.
                continue


def _unresolved_pages(parsed: ParsedDocument) -> set[int] | None:
    raw = parsed.raw_payload.get("unresolved_pages", [])
    if not isinstance(raw, list):
        return None
    pages = {item for item in raw if isinstance(item, int) and item > 0}
    return pages or None


def _merge_visual_results(
    parsed: ParsedDocument,
    visual: ParsedDocument,
) -> ParsedDocument:
    blocks = list(parsed.blocks)
    seen = {
        (block.page_number, block.text.strip().casefold(), block.block_type.casefold())
        for block in blocks
    }
    for block in visual.blocks:
        key = (block.page_number, block.text.strip().casefold(), block.block_type.casefold())
        if key not in seen:
            seen.add(key)
            blocks.append(block)
    blocks.sort(
        key=lambda block: (
            block.page_number or 0,
            block.bounding_box.top if block.bounding_box is not None else 0,
            block.text.casefold(),
        )
    )
    vision_pages = {block.page_number for block in visual.blocks if block.page_number is not None}
    unresolved = sorted((_unresolved_pages(parsed) or set()) - vision_pages)
    parser_chain = list(dict.fromkeys([*parsed.parser_chain, *visual.parser_chain]))
    return parsed.model_copy(
        update={
            "page_count": max(parsed.page_count, visual.page_count),
            "blocks": blocks,
            "language": parsed.language or visual.language,
            "detected_language": parsed.detected_language or visual.detected_language,
            "raw_payload": {
                **parsed.raw_payload,
                "unresolved_pages": unresolved,
                "vision": visual.raw_payload,
            },
            "warnings": [*parsed.warnings, *visual.warnings],
            "parser_chain": parser_chain,
            "vision_used": True,
            "low_confidence_blocks": [
                *parsed.low_confidence_blocks,
                *visual.low_confidence_blocks,
            ],
        }
    )


def _is_visual_input(path: Path, mime_type: str) -> bool:
    return (
        mime_type.startswith("image/")
        or path.suffix.casefold() == ".pdf"
        or mime_type == "application/pdf"
    )


def _is_pdf_input(path: Path, mime_type: str) -> bool:
    return path.suffix.casefold() == ".pdf" or mime_type == "application/pdf"


def _merge_block_lists(
    existing: list[ParsedBlock], additions: list[ParsedBlock]
) -> list[ParsedBlock]:
    """Merge parser output without duplicating identical page-level text."""

    blocks = list(existing)
    seen = {
        (block.page_number, block.text.strip().casefold(), block.block_type.casefold())
        for block in blocks
    }
    for block in additions:
        key = (block.page_number, block.text.strip().casefold(), block.block_type.casefold())
        if key not in seen:
            seen.add(key)
            blocks.append(block)
    return sorted(
        blocks,
        key=lambda block: (
            block.page_number or 0,
            block.bounding_box.top if block.bounding_box is not None else 0,
            block.text.casefold(),
        ),
    )


def _text_quality(blocks: list[ParsedBlock]) -> float:
    if not blocks:
        return 0.0
    text = "".join(block.text for block in blocks)
    if not text:
        return 0.0
    visible = sum(character.isalnum() for character in text)
    # Short text-only pages are valid (for example, a title slide), so use a
    # low threshold while still rejecting whitespace/garbled extraction.
    return min(1.0, visible / max(len(text), 1))


def _parsed_quality(parsed: ParsedDocument) -> float:
    scores = [block.confidence for block in parsed.blocks if block.confidence is not None]
    if scores:
        return sum(scores) / len(scores)
    return _text_quality(parsed.blocks)


def _ensure_parser_metadata(parsed: ParsedDocument, parser_name: str) -> ParsedDocument:
    chain = parsed.parser_chain or [parser_name]
    detected = parsed.detected_language or parsed.language
    return parsed.model_copy(
        update={
            "parser_chain": chain,
            "detected_language": detected,
            "ocr_used": parsed.ocr_used or "ocr" in parser_name.casefold(),
        }
    )


def _detect_language(blocks: list[ParsedBlock]) -> str | None:
    text = "".join(block.text for block in blocks)
    if not text:
        return None
    cjk = sum("\u4e00" <= character <= "\u9fff" for character in text)
    return "zh-CN" if cjk > max(2, len(text) // 20) else "en"


def _package_version(name: str) -> str:
    try:
        return version(name)
    except PackageNotFoundError:
        return "unknown"


def _docling_bounding_box(provenance: object) -> BoundingBox | None:
    raw = getattr(provenance, "bbox", None)
    if raw is None:
        return None

    def coordinate(*names: str) -> float | None:
        for name in names:
            value = getattr(raw, name, None)
            if isinstance(value, (int, float)):
                return float(value)
        return None

    left = coordinate("l", "left", "x0")
    top = coordinate("t", "top", "y0")
    right = coordinate("r", "right", "x1")
    bottom = coordinate("b", "bottom", "y1")
    if None in {left, top, right, bottom}:
        return None
    assert left is not None and top is not None and right is not None and bottom is not None
    origin = str(getattr(raw, "coord_origin", "page"))
    return BoundingBox(
        left=min(left, right),
        top=min(top, bottom),
        right=max(left, right),
        bottom=max(top, bottom),
        coordinate_space=f"docling:{origin}",
    )
